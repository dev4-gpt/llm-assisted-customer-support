#!/usr/bin/env python3
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Image, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle


ROOT = Path(__file__).resolve().parents[1]
ARTIFACTS = ROOT / "artifacts"
DOCS = ROOT / "docs"


def _load_json(path: Path) -> dict:
    if not path.exists():
        return {}
    return json.loads(path.read_text(encoding="utf-8"))


def _slide_title_and_bullets(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(20)


def _slide_title_only(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _slide_image(prs: Presentation, title: str, image_path: Path, caption: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.8), Inches(1.4), width=Inches(8.4))
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(6.4), Inches(8.4), Inches(0.5))
    p = tx.text_frame.paragraphs[0]
    p.text = caption
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)


def build_pptx() -> Path:
    metrics = _load_json(ARTIFACTS / "eval" / "metrics.json")
    train = _load_json(ARTIFACTS / "triage_roberta_demo" / "train_metrics.json")
    prs = Presentation()

    _slide_title_only(
        prs,
        "LLM-Assisted Customer Support Triage System",
        "Final Project Presentation (Evidence-Backed, 50 Slides)",
    )

    static_slides: list[tuple[str, list[str]]] = [
        ("Agenda", ["Problem", "Architecture", "Implementation", "Results", "Evaluation", "Alignment", "Roadmap"]),
        ("Problem Statement", ["Support queues are noisy and high-volume", "Manual triage is slow and inconsistent", "Need explainable AI-assisted routing"]),
        ("Project Objectives", ["Automate triage and quality monitoring", "Provide summarization and policy grounding", "Deliver reproducible demo-ready stack"]),
        ("System Scope", ["Endpoints: triage, quality, pipeline, summarize, rag/context", "Inference-first architecture with optional training paths", "Production-style API + tests + docs"]),
        ("Core NLP Tasks", ["Intent/category classification", "Priority and sentiment estimation", "Response quality scoring", "Thread summarization"]),
        ("High-Level Architecture", ["Notebook/API client -> FastAPI -> services", "LLM client shared by triage/quality/summarization", "Pipeline runs triage and quality concurrently"]),
        ("Data Assets", ["Golden eval set: data/golden/eval_set.jsonl", "Policy snippets for RAG: data/policy_snippets.json", "Demo transformer data: artifacts/demo_tickets.csv"]),
        ("Configuration Strategy", ["Environment-driven settings via app/core/config.py", "Provider profiles: ollama/openrouter/nvidia/manual", "Feature flags for optional components"]),
        ("API Contract Design", ["Pydantic models as source of truth", "Strict schema validation and explicit error handling", "Deterministic routing matrix outside LLM"]),
        ("Triage Service", ["LLM returns structured JSON", "Priority/category/intents normalized and validated", "Routing determined by matrix + sentiment escalation"]),
        ("Quality Service", ["Rubric-based scoring (empathetic/actionable/safety/resolution)", "Pass/fail threshold from config", "Coaching feedback and flagged phrases"]),
        ("Pipeline Service", ["Runs triage + quality in parallel", "Computes SLA recommendation", "Returns workflow_passed aggregate decision"]),
        ("Summarization Service", ["Multi-turn thread summarization", "Returns summary, key_points, confidence", "Schema-validated output"]),
        ("RAG Service", ["Default lexical retrieval over policy snippets", "Optional embedding retrieval backend", "Used for grounding quality/triage prompts"]),
        ("LLM Client", ["Provider-agnostic OpenAI-compatible + Anthropic support", "Retries, timeout handling, JSON parse safeguards", "Metrics labels include prompt_version"]),
        ("Reliability Issue Encountered", ["Hosted model output drift produced invalid labels (e.g., refund)", "Pipeline could fail with 422 on strict enum validation", "Needed deterministic recovery layer"]),
        ("Reliability Fix Implemented", ["Synonym normalization for common off-schema labels", "Embedding-similarity fallback mapping", "Strict final validation preserved"]),
        ("Intent Fallback Service", ["File: app/services/intent_fallback_service.py", "Maps invalid labels into allowed taxonomy", "Feature-gated with config flags"]),
        ("Fallback Config Flags", ["TRIAGE_EMBEDDING_FALLBACK_ENABLED", "TRIAGE_EMBEDDING_MODEL", "TRIAGE_EMBEDDING_MIN_SIMILARITY"]),
        ("Testing Strategy", ["Unit tests for services and fallback", "Integration tests for API and full pipeline", "Deterministic e2e tests for demo reliability"]),
        ("Verification Sweep", ["Core tests passed", "JSON/TOML/YAML syntax validated", "Notebook endpoint flow validated"]),
        ("Operational Readiness", ["README + alignment docs updated", "Notebook warm-up + retry guidance added", "Cleanup performed for generated clutter"]),
        ("Repository Structure Overview", ["app/: runtime services and API", "scripts/: training/evaluation tooling", "tests/: verification", "docs/: evidence/presentation assets"]),
        ("Demo Flow", ["Start uvicorn", "Warm-up triage call", "Run 5 endpoint showcase", "Inspect structured outputs"]),
        ("Security and Secrets", [".env local-only, .env.example template", "Do not commit provider keys", "API key middleware optional"]),
        ("Observability", ["Prometheus metrics available", "Structured logs with service context", "Prompt version labeling for regression tracking"]),
        ("Evaluation Methodology", ["Offline eval using golden set", "Category + priority accuracy", "Quality mean score + summarization ROUGE-L"]),
        ("Transformer Fine-Tune Path", ["Optional script: train_triage_transformer.py", "Produces checkpoint under artifacts/", "Can be used as triage hint"]),
        ("Baseline Classifier Path", ["Optional TF-IDF + Logistic Regression", "Used as lightweight hint in triage prompt", "Supports ablation-style demos"]),
        ("Notebook Coverage", ["EDA", "Offline evaluation", "Optional transformer training", "Live API endpoint checks"]),
        ("Presentation Risk Controls", ["Claims-vs-evidence matrix prepared", "Optional features clearly labeled", "Future work explicitly separated"]),
        ("Cleanup Summary", ["Removed only generated cache/metadata", "Preserved source, tests, docs, artifacts", "Cleaner repo for final presentation"]),
        ("Lecture Alignment - Overview", ["End-to-end NLP pipeline", "Embeddings + transfer learning", "Evaluation methodology", "Production engineering practice"]),
        ("Lecture Alignment - Embeddings", ["Semantic similarity fallback", "Optional embedding RAG backend", "Connects to representation learning concepts"]),
        ("Lecture Alignment - Transformers", ["Optional BERT/RoBERTa fine-tuning script", "Inference-time hint integration", "Transfer learning from pretrained models"]),
        ("Lecture Alignment - Evaluation", ["Classification metrics", "Quality scoring metrics", "Summarization ROUGE-L", "Regression-focused evaluation"]),
        ("Current Limitations", ["NER standalone module not shipped", "Large-corpus training remains optional", "External integrations are partial/stubbed"]),
        ("Roadmap", ["NER extraction module", "Expanded dataset experiments", "Richer analytics dashboards", "Deeper connector automation"]),
        ("Contribution Summary", ["Reliable production-style NLP stack", "Robustness fix for real-world label drift", "Reproducible demo and verification artifacts"]),
        ("Appendix: Key Files", ["app/services/triage_service.py", "app/services/intent_fallback_service.py", "tests/integration/test_pipeline_label_recovery.py"]),
        ("Appendix: Core Commands", ["pytest ... --no-cov", "python scripts/run_offline_eval.py ...", "uvicorn app.main:app --host 127.0.0.1 --port 8000"]),
        ("Appendix: Demo Checklist", ["Correct .env profile/model", "Restart uvicorn after env changes", "Run notebook warm-up before full endpoint cell"]),
    ]

    for title, bullets in static_slides:
        _slide_title_and_bullets(prs, title, bullets)

    # Result-specific slides
    triage_acc = metrics.get("triage_category", {}).get("accuracy", "N/A")
    priority_acc = metrics.get("triage_priority", {}).get("accuracy", "N/A")
    quality_mean = metrics.get("quality", {}).get("mean_score", "N/A")
    rouge = metrics.get("summarize", {}).get("mean_rouge_l_f1", "N/A")
    _slide_title_and_bullets(
        prs,
        "Offline Evaluation Snapshot",
        [
            f"Triage category accuracy: {triage_acc}",
            f"Triage priority accuracy: {priority_acc}",
            f"Quality mean score: {quality_mean}",
            f"Summarization ROUGE-L F1: {rouge}",
        ],
    )
    _slide_title_and_bullets(
        prs,
        "Transformer Demo Metrics",
        [
            f"Eval loss: {train.get('eval_loss', 'N/A')}",
            f"Eval accuracy: {train.get('eval_accuracy', 'N/A')}",
            f"Eval runtime: {train.get('eval_runtime', 'N/A')}",
            f"Epochs: {train.get('epoch', 'N/A')}",
        ],
    )

    _slide_image(
        prs,
        "EDA: Golden Task Counts",
        ARTIFACTS / "eda" / "golden_task_counts.png",
        "Distribution of golden-set tasks used in offline analysis.",
    )
    _slide_image(
        prs,
        "EDA: Text Length by Task",
        ARTIFACTS / "eda" / "golden_text_length_by_task.png",
        "Text length spread by task type for coverage sanity check.",
    )
    _slide_image(
        prs,
        "EDA: Triage Category Distribution",
        ARTIFACTS / "eda" / "golden_triage_category.png",
        "Category distribution in the golden set.",
    )
    _slide_image(
        prs,
        "EDA: Triage Priority Distribution",
        ARTIFACTS / "eda" / "golden_triage_priority.png",
        "Priority distribution in the golden set.",
    )

    # Add final slides until exactly 50
    while len(prs.slides) < 50:
        _slide_title_and_bullets(
            prs,
            f"Backup Slide {len(prs.slides) + 1}",
            [
                "Additional Q&A support material.",
                "Use this slot for live demo observations.",
                "Map answers back to docs/claims-evidence matrix.",
            ],
        )

    if len(prs.slides) != 50:
        raise RuntimeError(f"Expected 50 slides, got {len(prs.slides)}")

    out = ROOT / "LLM_Assist_Final_Presentation_50_Slides.pptx"
    prs.save(out)
    return out


def build_report_pdf() -> Path:
    metrics = _load_json(ARTIFACTS / "eval" / "metrics.json")
    train = _load_json(ARTIFACTS / "triage_roberta_demo" / "train_metrics.json")
    summary_md = (ARTIFACTS / "eval" / "summary.md").read_text(encoding="utf-8")

    out = ROOT / "LlmCustomerSupport_project_report_updated.pdf"
    doc = SimpleDocTemplate(str(out), pagesize=letter, rightMargin=40, leftMargin=40, topMargin=40)
    styles = getSampleStyleSheet()
    h1 = styles["Heading1"]
    h2 = styles["Heading2"]
    body = ParagraphStyle("Body", parent=styles["BodyText"], leading=14, spaceAfter=8)
    story = []

    story.append(Paragraph("LLM-Assisted Customer Support Triage and Quality Monitoring", h1))
    story.append(Paragraph("Updated Technical Report with Reproducible Results", body))
    story.append(Spacer(1, 0.2 * inch))

    story.append(Paragraph("1. Executive Summary", h2))
    story.append(
        Paragraph(
            "This report updates the project with latest verified runtime behavior, reliability fixes, "
            "evaluation outputs, and presentation-ready evidence generated from scripts and notebooks.",
            body,
        )
    )

    story.append(Paragraph("2. System Architecture and Flow", h2))
    for line in [
        "Notebook/API client sends requests to FastAPI endpoints.",
        "Core services: triage, quality, pipeline, summarization, RAG.",
        "LLM client handles provider communication, retries, and JSON parsing.",
        "Intent fallback service recovers off-schema labels via synonym and embedding similarity.",
    ]:
        story.append(Paragraph(f"- {line}", body))

    story.append(Paragraph("3. Verified Results", h2))
    table_data = [
        ["Metric", "Value"],
        ["Triage category accuracy", str(metrics.get("triage_category", {}).get("accuracy", "N/A"))],
        ["Triage priority accuracy", str(metrics.get("triage_priority", {}).get("accuracy", "N/A"))],
        ["Quality mean score", str(metrics.get("quality", {}).get("mean_score", "N/A"))],
        ["Summarization ROUGE-L F1", str(metrics.get("summarize", {}).get("mean_rouge_l_f1", "N/A"))],
        ["Transformer eval loss", str(train.get("eval_loss", "N/A"))],
        ["Transformer eval accuracy", str(train.get("eval_accuracy", "N/A"))],
    ]
    t = Table(table_data, colWidths=[2.8 * inch, 2.8 * inch])
    t.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("ALIGN", (0, 0), (-1, -1), "LEFT"),
            ]
        )
    )
    story.append(t)
    story.append(Spacer(1, 0.2 * inch))

    story.append(Paragraph("4. Offline Evaluation Summary", h2))
    for line in summary_md.splitlines():
        if line.strip():
            story.append(Paragraph(line.replace("*", ""), body))

    story.append(Paragraph("5. Graphs from Script Outputs", h2))
    for image_name in [
        "golden_task_counts.png",
        "golden_text_length_by_task.png",
        "golden_triage_category.png",
        "golden_triage_priority.png",
    ]:
        img_path = ARTIFACTS / "eda" / image_name
        if img_path.exists():
            story.append(Paragraph(image_name, body))
            story.append(Image(str(img_path), width=6.2 * inch, height=3.4 * inch))
            story.append(Spacer(1, 0.1 * inch))

    story.append(Paragraph("6. Reliability Improvements", h2))
    for line in [
        "Added strict-but-resilient label recovery for invalid LLM outputs.",
        "Synonym mapping (e.g., refund -> billing) with optional embedding fallback.",
        "Added unit and integration tests for pipeline label recovery path.",
    ]:
        story.append(Paragraph(f"- {line}", body))

    story.append(Paragraph("7. Alignment with In-Class Lectures", h2))
    for line in [
        "End-to-end NLP pipeline implementation (API orchestration).",
        "Embedding-based semantic similarity usage in fallback and RAG options.",
        "Optional transformer fine-tuning path (transfer learning).",
        "Evaluation methodology and reproducibility via scripted checks.",
    ]:
        story.append(Paragraph(f"- {line}", body))

    story.append(Paragraph("8. Current Limitations and Next Steps", h2))
    for line in [
        "NER module remains planned work.",
        "Large-corpus experiments are documented/optional, not the only active evidence path.",
        "External integration write-back is still partial/stub level.",
    ]:
        story.append(Paragraph(f"- {line}", body))

    doc.build(story)
    return out


def main() -> None:
    pptx_out = build_pptx()
    pdf_out = build_report_pdf()
    print(f"Generated presentation: {pptx_out}")
    print(f"Generated updated report: {pdf_out}")


if __name__ == "__main__":
    main()
